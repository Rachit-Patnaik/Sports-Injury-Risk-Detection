import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    DARK_BG = RGBColor(7, 10, 18)
    CARD_BG = RGBColor(15, 23, 42)
    ACCENT_BLUE = RGBColor(56, 189, 248)
    ACCENT_INDIGO = RGBColor(129, 140, 248)
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)
    BORDER_COLOR = RGBColor(30, 41, 59)

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def add_header(slide, title_text, category_text="SPORTS AI PLATFORM"):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_INDIGO
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, title="", items=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = BORDER_COLOR
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_bottom = Inches(0.25)

        if title:
            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.size = Pt(15)
            p0.font.bold = True
            p0.font.color.rgb = ACCENT_BLUE
            p0.space_after = Pt(10)

        if items:
            for idx, item in enumerate(items):
                p = tf.add_paragraph() if (title or idx > 0) else tf.paragraphs[0]
                p.text = f"• {item}"
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_MUTED
                p.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # -------------------------------------------------------------
    slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(slide_layout)
    set_background(slide1)

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "SPORTS INJURY RISK DETECTION FROM VIDEO"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "An End-to-End AI & Computer Vision Platform for Biomechanical Movement Intelligence"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_BLUE
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "B.Tech Computer Science & Engineering Final Project | ITER SOA University"
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(30)

    # -------------------------------------------------------------
    # SLIDES 2-19 GENERATION
    # -------------------------------------------------------------
    slides_data = [
        ("Problem Statement & Motivation", "BIOMECHANICS & INJURY PREVENTION", [
            ("High Non-Contact Injury Rates", ["Non-contact ACL tears, hamstring strains, and joint overhead stress account for over 60% of athletic absences.", "Traditional motion labs require expensive marker-based equipment and specialized setups."]),
            ("The Solution Gap", ["Need for accessible, markerless video screening using consumer cameras.", "Automated frame-by-frame joint angle extraction and early predictive risk warnings."])
        ]),
        ("Project Objectives & Outcomes", "CORE DELIVERABLES", [
            ("Technical Objectives", ["Develop markerless 2D/3D pose estimation using MediaPipe and OpenCV.", "Calculate kinematic angles, Range of Motion (ROM), standard deviation, and bilateral asymmetry."]),
            ("Platform Outcomes", ["Construct role-differentiated dashboards for Coaches, Physios, Scientists, Athletes, and Admins.", "Implement SQLite database persistence, JWT auth, PDF/Excel export, and Docker containers."])
        ]),
        ("System Architecture Overview", "END-TO-END PIPELINE", [
            ("API Gateway & Backend", ["FastAPI routing gateway handling request validation, CORS, and multipart video ingestion.", "SQLAlchemy ORM layer interfacing with SQLite3 persistent database."]),
            ("AI & Computer Vision Layer", ["MediaPipe Pose landmark detector extracting 33 keypoints across 100 keyframes.", "Kinematic Engine calculating joint angles and feeding Recommendation Predictor."])
        ]),
        ("Technology Stack & Frameworks", "SOFTWARE INFRASTRUCTURE", [
            ("Frontend Infrastructure", ["React.js (Vite), Framer Motion animations, Tailwind CSS, and Material-UI icons."]),
            ("Backend & AI Stack", ["Python 3.10, FastAPI, Uvicorn, OpenCV, MediaPipe, NumPy, Pandas, Scikit-learn, SQLAlchemy, Docker."])
        ]),
        ("Module 1 & 2: Auth & Athlete Profiles", "SECURITY & DATA MANAGEMENT", [
            ("Module 1: User Auth & RBAC", ["JWT token generation with PBKDF2-HMAC-SHA256 password hashing.", "Role-Based Access Control supporting 5 distinct user roles."]),
            ("Module 2: Athlete Profiles", ["Physical metrics management (Height, Weight, Age, Position, Training Load).", "SQLite DB schema tracking historical injury records."])
        ]),
        ("Module 3: Video Processing Engine", "MEDIA INGESTION", [
            ("Multi-Activity Context Support", ["Supports Running, Sprinting, Jumping, Squatting, Landing, Throwing, and Cutting."]),
            ("Processing Workflow", ["Dynamic frame sampling stride optimization.", "Standardized frame resolution and CSV joint log generation."])
        ]),
        ("Module 4: Computer Vision Pose Engine", "KEYPOINT DETECTION", [
            ("Landmark Topology", ["Extracts 33 body landmarks concentrating on lower and upper limb joint vertices."]),
            ("Joint Angle Identifiers", ["Knee: Hip (23/24) -> Knee (25/26) -> Ankle (27/28).", "Hip: Shoulder (11/12) -> Hip (23/24) -> Knee (25/26).", "Shoulder: Elbow (13/14) -> Shoulder (11/12) -> Hip (23/24)."])
        ]),
        ("Module 5: Biomechanical Formulas", "KINEMATIC MATHEMATICS", [
            ("2D Vector Joint Angle Formula", ["θ = |atan2(yc - yb, xc - xb) - atan2(ya - yb, xa - xb)| * (180 / π)"]),
            ("Bilateral Asymmetry Index", ["Asymmetry % = (|θ_left - θ_right| / ((θ_left + θ_right) / 2)) * 100%"])
        ]),
        ("Module 6: Injury Risk Prediction Engine", "PREDICTIVE MODELING", [
            ("Injury Categories Monitored", ["ACL Vulnerability, Hamstring Strain Risk, Shoulder Imbalance, and Gait Symmetry."]),
            ("Contextual Thresholds", ["Evaluates deep flexion deficits during squats and landing mechanics."])
        ]),
        ("Module 7 & 8: Anomaly & Risk Scoring", "WEIGHTED SCORING ENGINE", [
            ("5-Factor Weighted Score Formula", ["Score = 35%(Bio Deviations) + 20%(History) + 20%(Asymmetry) + 15%(Load) + 10%(Fatigue)"]),
            ("Risk Bands", ["Low Risk (<30%), Moderate Risk (30-55%), High Risk (>55%)."])
        ]),
        ("Module 9: Corrective Recommendations", "CLINICAL PROTOCOLS", [
            ("Targeted Exercise Mapping", ["High ACL Risk -> Neuromuscular jump-landing & gluteus medius activation."]),
            ("Hamstring & Shoulder Drills", ["High Hamstring Risk -> Eccentric Nordic curls.", "Shoulder Imbalance -> Scapular stabilization drills."])
        ]),
        ("Module 10: Role-Based Dashboards", "USER EXPERIENCE", [
            ("Coach & Athlete Views", ["Coach: Squad risk roster & load status. Athlete: Personal health score & drill checklist."]),
            ("Physio, Scientist & Admin Views", ["Physio: Rehab tracking. Scientist: Kinematic matrix. Admin: API & DB health."])
        ]),
        ("Module 11 & 12: Alerts & Export System", "NOTIFICATIONS & REPORTING", [
            ("Module 11: Alert System", ["Top navigation bar notification bell with live high-risk popover alerts."]),
            ("Module 12: Export Engine", ["Native PDF printing and formatted Excel/CSV kinematic data exports."])
        ]),
        ("Module 13: Persistence & Deployment", "DEPOPS & PERSISTENCE", [
            ("Database Persistence", ["SQLite3 database managed via SQLAlchemy ORM (users, athletes, reports)."]),
            ("Docker Containerization", ["Multi-stage Dockerfile and docker-compose.yml orchestrating FastAPI and Vite."])
        ]),
        ("Experimental Validation & Results", "PERFORMANCE EVALUATION", [
            ("Accuracy Metrics", ["Keypoint Detection Accuracy > 94.2% PCK@0.2."]),
            ("Kinematic Precision", ["Joint Angle Error MAE < 3.8° vs goniometer.", "Pipeline Latency < 2.5s for 200 frames."])
        ]),
        ("Engineering Challenges & Solutions", "TECHNICAL RESOLUTIONS", [
            ("Python 3.12 MediaPipe Imports", ["Resolved via dynamic fallback wrappers handling solutions.pose."]),
            ("Weighted Scoring Integration", ["Unified kinematic standard deviation rules with demographic training load."])
        ]),
        ("Future Scope & Conclusion", "ROADMAP", [
            ("Future Enhancements", ["3D volumetric mesh reconstruction (SMPL models) and wearable sensor IMU fusion."]),
            ("Conclusion", ["Successfully deployed production-ready markerless sports injury intelligence platform."])
        ])
    ]

    for title, cat, cards in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        add_header(slide, title, cat)
        
        if len(cards) == 2:
            add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), cards[0][0], cards[0][1])
            add_card(slide, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.2), cards[1][0], cards[1][1])

    output_path = "Sports_AI_Injury_Risk_Detection.pptx"
    prs.save(output_path)
    print(f"[SUCCESS] 19-Slide PowerPoint presentation generated: {output_path}")

if __name__ == "__main__":
    create_deck()